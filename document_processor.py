"""
document_processor.py — Extraction de texte depuis divers formats (PDF, DOCX, ODT, PPTX) et chunking intelligent.
"""

import re
import io
from dataclasses import dataclass, field
from typing import List, Literal, BinaryIO, Any

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from odf.opendocument import load as load_odf
from odf import text, draw, table
from odf.teletype import extractText

from llm_service import count_tokens, _encoder


@dataclass
class TextChunk:
    """Un morceau de texte extrait du document avec ses métadonnées."""
    text: str
    source_pages: List[int] = field(default_factory=list)
    token_count: int = 0
    source_document: str = ""


def _extract_from_pdf(file: BinaryIO) -> List[dict]:
    """Extrait le texte d'un PDF page par page."""
    pages = []
    try:
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages):
                text_content = page.extract_text() or ""
                # Nettoyage basique
                text_content = re.sub(r'\s+', ' ', text_content).strip()
                text_content = re.sub(r'(\n\s*){3,}', '\n\n', text_content)
                if text_content:
                    pages.append({"page": i + 1, "text": text_content})
    except Exception as e:
        print(f"Erreur lors de l'extraction PDF : {e}")
    return pages


def _extract_from_docx(file: BinaryIO) -> List[dict]:
    """Extrait le texte d'un fichier DOCX en gérant les sauts de page."""
    pages = []
    try:
        doc = DocxDocument(file)
        
        current_page_text = []
        page_num = 1
        
        # namespaces pour XPath
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

        def flush_page():
            nonlocal current_page_text, page_num
            text_content = "\n\n".join(current_page_text).strip()
            if text_content:
                pages.append({"page": page_num, "text": text_content})
                page_num += 1
            current_page_text = []

        for para in doc.paragraphs:
            # Vérifier si le paragraphe a un saut de page avant (propriété de style)
            if para._element.xpath('.//w:pPr/w:pageBreakBefore', namespaces=ns):
                flush_page()

            para_parts = []
            
            for run in para.runs:
                # Vérifier les sauts de page à l'intérieur des runs
                # w:br type="page" (saut manuel) ou w:lastRenderedPageBreak (saut automatique Word)
                breaks = run._element.xpath('.//w:br[@w:type="page"] | .//w:lastRenderedPageBreak', namespaces=ns)
                
                if breaks:
                    # S'il y a un break, on finit ce qui précédait dans le run
                    # Note: on ne peut pas facilement splitter le texte au milieu d'un run 
                    # mais en général les breaks sont entre les textes ou dans des runs dédiés
                    
                    # On flush le texte accumulé jusque là dans le paragraphe
                    if run.text:
                        para_parts.append(run.text)
                    
                    if para_parts:
                        current_page_text.append(" ".join(para_parts))
                        para_parts = []
                    
                    flush_page()
                else:
                    if run.text:
                        para_parts.append(run.text)
            
            if para_parts:
                current_page_text.append(" ".join(para_parts))

        # Ne pas oublier la dernière page
        flush_page()

        # Fallback si rien n'a été extrait (ex: document vide ou formatage étrange)
        if not pages and doc.paragraphs:
            full_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            if full_text:
                pages.append({"page": 1, "text": "\n\n".join(full_text)})

    except Exception as e:
        print(f"Erreur lors de l'extraction DOCX : {e}")
    return pages


def _extract_from_pptx(file: BinaryIO) -> List[dict]:
    """Extrait le texte d'un fichier PPTX (slide par slide)."""
    pages = []
    try:
        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            text_content = "\n".join(slide_text)
            if text_content:
                pages.append({"page": i + 1, "text": text_content})
    except Exception as e:
        print(f"Erreur lors de l'extraction PPTX : {e}")
    return pages


def _extract_from_odf(file: BinaryIO) -> List[dict]:
    """
    Extrait le texte d'un fichier ODT / ODP / ODS.
    - ODT : tente de respecter les sauts de page (hard & soft)
    - ODP : une slide = une "page"
    - Retour : liste de {"page": int, "text": str}
    """
    pages = []
    try:
        doc = load_odf(file)
        file.seek(0)  # au cas où

        # ───────────────────────────────────────────────
        # 1. Cas présentation → un slide = une page
        # ───────────────────────────────────────────────
        slide_elements = doc.getElementsByType(draw.Page)
        if slide_elements:
            for i, slide in enumerate(slide_elements, 1):
                # On prend tout le texte de la slide (formes, notes, etc.)
                slide_text = extractText(slide).strip()
                slide_text = re.sub(r'\s{2,}', ' ', slide_text)
                slide_text = re.sub(r'\n\s*\n+', '\n\n', slide_text)
                if slide_text:
                    pages.append({"page": i, "text": slide_text})
            if pages:
                return pages

        # ───────────────────────────────────────────────
        # 2. Cas tableur (ODS) → une feuille = une "page"
        # ───────────────────────────────────────────────
        tables = doc.getElementsByType(table.Table)
        if tables and not slide_elements:
            for i, tbl in enumerate(tables, 1):
                table_text = []
                for row in tbl.getElementsByType(table.TableRow):
                    row_cells = []
                    for cell in row.getElementsByType(table.TableCell):
                        cell_text = extractText(cell).strip()
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        table_text.append(" | ".join(row_cells))
                if table_text:
                    pages.append({
                        "page": i,
                        "text": "\n".join(table_text)
                    })
            if pages:
                return pages

        # ───────────────────────────────────────────────
        # 3. Cas texte (ODT) → on essaie de détecter les sauts de page
        # ───────────────────────────────────────────────
        current_page = []
        page_number = 1

        def flush_current_page():
            nonlocal page_number
            if current_page:
                text_block = "\n\n".join(current_page).strip()
                text_block = re.sub(r'\s{3,}', ' ', text_block)
                text_block = re.sub(r'\n\s*\n{2,}', '\n\n', text_block)
                if text_block:
                    pages.append({"page": page_number, "text": text_block})
                    page_number += 1
            current_page.clear()

        # On parcourt tous les éléments text:p et text:h
        for elem in doc.getElementsByType(text.P) + doc.getElementsByType(text.H):
            # ─── Détection saut de page AVANT ───────────────────────
            style_name = elem.getAttribute("stylename") or elem.getAttribute("style-name")
            if style_name:
                style = doc.getStyleByName(style_name)
                if style:
                    pp = style.getElementsByType(text.ParagraphProperties)
                    if pp:
                        br_before = pp[0].getAttribute("breakbefore") or pp[0].getAttribute("page-break-before")
                        if br_before and br_before.lower() in ("page", "always"):
                            flush_current_page()

            # ─── Extraction du texte du paragraphe ──────────────────
            para_text = extractText(elem).strip()
            if not para_text:
                continue

            # On regarde s'il y a des soft-page-breaks à l'intérieur
            # (c'est rare dans les <text:p>, mais possible dans les spans)
            children = elem.childNodes
            has_soft_break = any(
                isinstance(n, text.SoftPageBreak) or
                (hasattr(n, "tagName") and "soft-page-break" in n.tagName.lower())
                for n in children if hasattr(n, "tagName")
            )

            if has_soft_break:
                # Approximation : on flush avant le soft break
                current_page.append(para_text)
                flush_current_page()
            else:
                current_page.append(para_text)

            # ─── Détection saut de page APRÈS ───────────────────────
            if style_name and style:
                pp = style.getElementsByType(text.ParagraphProperties)
                if pp:
                    br_after = pp[0].getAttribute("breakafter") or pp[0].getAttribute("page-break-after")
                    if br_after and br_after.lower() in ("page", "always"):
                        flush_current_page()

        # Ne pas oublier la dernière page
        flush_current_page()

        # Fallback ultra-simple si on n'a rien trouvé
        if not pages:
            all_paragraphs = []
            for elem in doc.getElementsByType(text.P) + doc.getElementsByType(text.H):
                t = extractText(elem).strip()
                if t:
                    all_paragraphs.append(t)
            if all_paragraphs:
                pages.append({
                    "page": 1,
                    "text": "\n\n".join(all_paragraphs)
                })

    except Exception as e:
        print(f"Erreur extraction ODF/ODT : {e}")
        # Fallback ultime : texte brut
        try:
            raw_text = extractText(doc).strip()
            if raw_text:
                pages = [{"page": 1, "text": raw_text}]
        except:
            pass

    return pages
    
def _extract_from_txt(file: BinaryIO) -> List[dict]:
    """Extrait le texte d'un fichier texte simple."""
    try:
        content = file.read().decode("utf-8", errors="replace").strip()
        return [{"page": 1, "text": content}] if content else []
    except Exception as e:
        print(f"Erreur lors de l'extraction TXT : {e}")
        return []


def extract_text_from_file(file: BinaryIO) -> List[dict]:
    """
    Extrait le texte depuis un fichier selon son extension.
    Détecte automatiquement le type via file.name.
    """
    if not hasattr(file, "name"):
        # Fallback si pas de nom, essayons PDF par défaut ou erreur
        return _extract_from_pdf(file)

    filename = file.name.lower()
    
    # Reset le pointeur du fichier au début au cas où
    file.seek(0)
    
    if filename.endswith(".pdf"):
        return _extract_from_pdf(file)
    elif filename.endswith(".docx"):
        return _extract_from_docx(file)
    elif filename.endswith(".pptx"):
        return _extract_from_pptx(file)
    elif filename.endswith((".odt", ".odp", ".ods")):
        return _extract_from_odf(file)
    elif filename.endswith(".txt"):
        return _extract_from_txt(file)
    else:
        # Tenter PDF par défaut si extension inconnue mais accepté
        return _extract_from_pdf(file)


def split_into_paragraphs(pages: List[dict]) -> List[TextChunk]:
    """
    Sépare le texte en paragraphes.
    """
    paragraphs = []
    for page_data in pages:
        page_num = page_data["page"]
        text_content = page_data["text"]
        parts = re.split(r'\n\n+', text_content)
        for part in parts:
            part = part.strip()
            if part and len(part) > 20:
                paragraphs.append(TextChunk(
                    text=part,
                    source_pages=[page_num],
                    token_count=count_tokens(part)
                ))
    return paragraphs


def chunk_text(
    pages: List[dict],
    max_tokens: int = 2000,
    overlap_tokens: int = 200
) -> List[TextChunk]:
    """
    Découpe le texte complet en chunks de taille max_tokens avec overlap.
    """
    full_text = ""
    page_spans = []
    
    for page_data in pages:
        header = f"\n\n[Début Page {page_data['page']}]\n"
        footer = f"\n[Fin Page {page_data['page']}]"
        
        page_content = header + page_data["text"] + footer
        
        start_idx = len(full_text)
        full_text += page_content
        end_idx = len(full_text)
        
        page_spans.append((start_idx, end_idx, page_data["page"]))

    tokens = _encoder.encode(full_text)
    total_tokens = len(tokens)
    
    if total_tokens == 0:
        return []

    chunks = []
    start_token = 0
    
    while start_token < total_tokens:
        end_token = min(start_token + max_tokens, total_tokens)
        chunk_tokens_list = tokens[start_token:end_token]
        chunk_text_str = _encoder.decode(chunk_tokens_list)
        
        prefix_text = _encoder.decode(tokens[:start_token])
        chunk_char_start = len(prefix_text)
        chunk_char_end = chunk_char_start + len(chunk_text_str)
        
        source_pages = []
        for p_start, p_end, p_num in page_spans:
            if max(chunk_char_start, p_start) < min(chunk_char_end, p_end):
                if p_num not in source_pages:
                    source_pages.append(p_num)
        
        chunks.append(TextChunk(
            text=chunk_text_str.strip(),
            source_pages=source_pages,
            token_count=len(chunk_tokens_list)
        ))
        
        if end_token >= total_tokens:
            break
        start_token = end_token - overlap_tokens

    return chunks


def split_into_pages(pages: List[dict]) -> List[TextChunk]:
    """
    Sépare le texte par page (ou section logique selon format).
    """
    chunks = []
    for page_data in pages:
        text_content = page_data["text"].strip()
        if text_content:
            chunks.append(TextChunk(
                text=text_content,
                source_pages=[page_data["page"]],
                token_count=count_tokens(text_content)
            ))
    return chunks


def extract_and_chunk(
    file: BinaryIO,
    mode: Literal["page", "token"] = "page",
    max_tokens: int = 10000,
    overlap_tokens: int = 200
) -> List[TextChunk]:
    """
    Pipeline complet : extraction + chunking selon le mode choisi.
    """
    pages = extract_text_from_file(file)
    
    if not pages:
        return []

    if mode == "page":
        return split_into_pages(pages)
    
    elif mode == "token":
        return chunk_text(pages, max_tokens, overlap_tokens)
    
    else:
        raise ValueError(f"Mode inconnu : {mode}")


def get_full_text(pages: List[dict]) -> str:
    """Retourne le texte complet concaténé."""
    return "\n\n".join(p["text"] for p in pages)


def get_text_stats(file: BinaryIO) -> dict:
    """Retourne des statistiques sur le document."""
    pages = extract_text_from_file(file)
    full_text = get_full_text(pages)
    total_tokens = count_tokens(full_text)
    return {
        "num_pages": len(pages),
        "total_chars": len(full_text),
        "total_tokens": total_tokens,
        "avg_tokens_per_page": total_tokens // max(len(pages), 1)
    }


def get_text_stats_multiple(files: List[BinaryIO]) -> dict:
    """Retourne des statistiques globales et par document pour plusieurs fichiers."""
    per_doc = []
    total_pages = 0
    total_chars = 0
    total_tokens = 0

    for file in files:
        file.seek(0)
        stats = get_text_stats(file)
        file.seek(0)
        stats["name"] = getattr(file, "name", "inconnu")
        per_doc.append(stats)
        total_pages += stats["num_pages"]
        total_chars += stats["total_chars"]
        total_tokens += stats["total_tokens"]

    return {
        "num_pages": total_pages,
        "total_chars": total_chars,
        "total_tokens": total_tokens,
        "avg_tokens_per_page": total_tokens // max(total_pages, 1),
        "num_documents": len(files),
        "per_document": per_doc,
    }


def extract_and_chunk_multiple(
    files: List[BinaryIO],
    mode: Literal["page", "token"] = "page",
    max_tokens: int = 10000,
    overlap_tokens: int = 200
) -> List[TextChunk]:
    """
    Pipeline complet multi-documents : extraction + chunking pour chaque fichier.
    Chaque TextChunk conserve le nom du document source.
    """
    all_chunks = []
    for file in files:
        file.seek(0)
        doc_name = getattr(file, "name", "inconnu")
        chunks = extract_and_chunk(file, mode=mode, max_tokens=max_tokens, overlap_tokens=overlap_tokens)
        for chunk in chunks:
            chunk.source_document = doc_name
        all_chunks.extend(chunks)
        file.seek(0)
    return all_chunks